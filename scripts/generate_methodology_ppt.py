"""Generate the Gold Forecast Methodology PPT.

Mirrors the conversational explanation given in chat — 5 levers,
worked example, 4 scenarios, how to read the dashboard. Output is
saved to web/public/ so the dashboard can offer it as a one-click
download.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parent.parent
OUT_PATH = ROOT / "web" / "public" / "gold-forecast-methodology.pptx"

# Brand palette — matches the dashboard
GOLD_900 = RGBColor(0x4A, 0x33, 0x06)
GOLD_700 = RGBColor(0x80, 0x57, 0x0F)
GOLD_500 = RGBColor(0xC9, 0x90, 0x25)
GOLD_300 = RGBColor(0xF0, 0xC2, 0x60)
GOLD_50  = RGBColor(0xFF, 0xF8, 0xE5)
CREAM    = RGBColor(0xFB, 0xF7, 0xEE)
INK      = RGBColor(0x1F, 0x1A, 0x14)
INK_2    = RGBColor(0x4A, 0x42, 0x36)
POS      = RGBColor(0x4F, 0x7F, 0x4E)
NEG      = RGBColor(0xC5, 0x4F, 0x4F)
BLUE     = RGBColor(0x4A, 0x90, 0xC5)
PURPLE   = RGBColor(0x8C, 0x5D, 0x9A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color: RGBColor) -> None:
    """Fill the slide background with a solid color."""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H,
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    slide.shapes._spTree.insert(2, rect._element)


def add_left_accent(slide, color: RGBColor = None) -> None:
    """Vertical gold bar on the left margin."""
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SLIDE_H,
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = color or GOLD_500
    accent.line.fill.background()


def add_textbox(slide, left, top, width, height, text, *,
                size=18, bold=False, color=INK,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                font="Helvetica"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullet_list(slide, left, top, width, height, items, *,
                    size=14, color=INK_2, bold_first_word=False,
                    line_spacing=1.4):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        bullet_run = p.add_run()
        bullet_run.text = "• "
        bullet_run.font.size = Pt(size)
        bullet_run.font.color.rgb = GOLD_700
        bullet_run.font.bold = True
        if isinstance(item, tuple):
            head, body = item
            r1 = p.add_run()
            r1.text = head
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = INK
            r1.font.name = "Helvetica"
            r2 = p.add_run()
            r2.text = " " + body
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
            r2.font.name = "Helvetica"
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.name = "Helvetica"
    return tb


def add_chip(slide, left, top, label: str, fill: RGBColor,
             fg: RGBColor = None, width=Inches(2.2), height=Inches(0.5),
             size=11, bold=True):
    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = fill
    chip.line.color.rgb = fill
    chip.adjustments[0] = 0.5
    tf = chip.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = fg or CREAM
    r.font.name = "Helvetica"


def add_table(slide, left, top, headers, rows, *,
              col_widths=None, row_height=Inches(0.45),
              header_fill=GOLD_700, header_fg=CREAM,
              row_fill_alt=GOLD_50, row_fg=INK, font_size=11,
              header_size=11):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    if col_widths is None:
        total = Inches(11)
        col_widths = [total // n_cols] * n_cols
    total_w = sum(col_widths, Emu(0))
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, total_w, row_height * n_rows)
    table = tbl_shape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    # Header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.margin_left = Inches(0.1)
        cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.06)
        cell.margin_bottom = Inches(0.06)
        tf = cell.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h
        r.font.size = Pt(header_size)
        r.font.bold = True
        r.font.color.rgb = header_fg
        r.font.name = "Helvetica"
    # Body
    for r_idx, row in enumerate(rows, start=1):
        fill = row_fill_alt if r_idx % 2 == 1 else CREAM
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            # Support (text, color) tuples for emphasis
            if isinstance(val, tuple):
                run.text = val[0]
                run.font.color.rgb = val[1]
                if len(val) > 2 and val[2]:
                    run.font.bold = True
            else:
                run.text = str(val)
                run.font.color.rgb = row_fg
            run.font.size = Pt(font_size)
            run.font.name = "Helvetica"
    return table


def add_footer(slide, page: int, total: int) -> None:
    add_textbox(
        slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
        "Gold Forecast Methodology  ·  goldorg dashboard",
        size=9, color=INK_2,
    )
    add_textbox(
        slide, Inches(11.5), Inches(7.05), Inches(1.4), Inches(0.3),
        f"{page} / {total}",
        size=9, color=INK_2, align=PP_ALIGN.RIGHT,
    )


# ────────────────────────────────────────────────────────────────────
# SLIDES
# ────────────────────────────────────────────────────────────────────
def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    slides_count_placeholder = 18

    # 1. COVER
    s = prs.slides.add_slide(blank)
    add_bg(s, GOLD_900)
    # Gold accent bar across the top
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.7), SLIDE_W, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD_500
    bar.line.fill.background()
    add_textbox(s, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.5),
                "GOLD FORECAST · METHODOLOGY",
                size=14, bold=True, color=GOLD_300, font="Helvetica")
    add_textbox(s, Inches(0.8), Inches(3.0), Inches(11.5), Inches(2.2),
                "How the dashboard predicts gold returns",
                size=44, bold=True, color=CREAM, font="Helvetica")
    add_textbox(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.2),
                "From first principles — what gold is, what moves it, "
                "and how to use the model.",
                size=18, color=GOLD_300)
    add_chip(s, Inches(0.8), Inches(6.0), "5 LEVERS", GOLD_500, INK)
    add_chip(s, Inches(3.1), Inches(6.0), "1 RECIPE", GOLD_500, INK)
    add_chip(s, Inches(5.4), Inches(6.0), "4 SCENARIOS", GOLD_500, INK)
    add_chip(s, Inches(7.7), Inches(6.0), "1 YOU", GOLD_500, INK)

    # 2. WHY DOES GOLD HAVE VALUE
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    add_left_accent(s)
    add_textbox(s, Inches(0.6), Inches(0.5), Inches(0.5), Inches(0.4),
                "01", size=18, bold=True, color=GOLD_500)
    add_textbox(s, Inches(0.6), Inches(0.9), Inches(11), Inches(0.8),
                "Why does gold have value at all?",
                size=32, bold=True, color=INK)
    add_textbox(s, Inches(0.6), Inches(1.8), Inches(11), Inches(0.6),
                "It doesn't pay interest. It doesn't earn profits. It just sits there shiny.",
                size=14, color=INK_2)
    add_textbox(s, Inches(0.6), Inches(2.4), Inches(11), Inches(0.5),
                "So why does any money flow into it? Three reasons, every single year:",
                size=14, color=INK_2)
    add_table(s, Inches(0.6), Inches(3.1),
              headers=["BUYER", "WHY THEY BUY", "SHARE"],
              rows=[
                  ["Wedding & jewellery", "India + China + Middle East · cultural, dowries, festivals", "~45%"],
                  ["Investors saving wealth", "You, hedge funds, ETFs · protection from inflation + crisis", "~30%"],
                  ["Central banks", "China + India + Russia + Turkey · less USD dependency", "~20%"],
                  ["Industry", "Chips, dentistry · actually using the metal", "~5%"],
              ],
              col_widths=[Inches(2.4), Inches(7.4), Inches(1.2)])
    add_textbox(s, Inches(0.6), Inches(5.7), Inches(11), Inches(0.5),
                "Mines pull ~3,500 tonnes out of the ground per year. Recycling adds ~1,200 tonnes.",
                size=12, color=INK_2)
    add_textbox(s, Inches(0.6), Inches(6.15), Inches(11), Inches(0.5),
                "Demand > Supply → price UP.  Supply > Demand → price DOWN.",
                size=13, bold=True, color=GOLD_700)
    add_footer(s, 2, slides_count_placeholder)

    # 3. WHY DOES PRICE CHANGE MONTH TO MONTH
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    add_left_accent(s)
    add_textbox(s, Inches(0.6), Inches(0.5), Inches(0.5), Inches(0.4),
                "02", size=18, bold=True, color=GOLD_500)
    add_textbox(s, Inches(0.6), Inches(0.9), Inches(11), Inches(0.8),
                "Why does the PRICE change month-to-month?",
                size=32, bold=True, color=INK)
    add_textbox(s, Inches(0.6), Inches(1.9), Inches(11), Inches(0.6),
                "Wedding demand is steady. Mine output is steady. So why does gold swing 20-30% a year?",
                size=14, color=INK_2)
    add_textbox(s, Inches(0.6), Inches(2.6), Inches(11), Inches(0.6),
                "Answer: big investors.",
                size=22, bold=True, color=GOLD_700)
    add_textbox(s, Inches(0.6), Inches(3.4), Inches(12), Inches(2.5),
                "When a hedge fund has $1 billion to allocate, they're asking every day:\n\n"
                "  → US bonds (paying 4% safely)\n"
                "  → Gold (paying 0% but protects in a crash)\n"
                "  → Stocks (risky)\n"
                "  → Cash (inflation eats it)\n\n"
                "Their answer changes when 5 things change. Those 5 things are this model.",
                size=15, color=INK_2)
    add_footer(s, 3, slides_count_placeholder)

    # 4. THE 5 LEVERS OVERVIEW
    s = prs.slides.add_slide(blank)
    add_bg(s, GOLD_900)
    add_textbox(s, Inches(0.6), Inches(0.5), Inches(11), Inches(0.5),
                "THE 5 LEVERS", size=14, bold=True, color=GOLD_300)
    add_textbox(s, Inches(0.6), Inches(1.0), Inches(11), Inches(0.8),
                "The whole show, in 5 things",
                size=32, bold=True, color=CREAM)
    # 5 columns
    levers = [
        ("US 10y Yield",        "Bonds compete with gold for safe-haven flows",   GOLD_500),
        ("US Debt/GDP",         "Fiscal stress → gold as escape valve",            BLUE),
        ("US CPI Inflation",    "Gold preserves purchasing power",                 NEG),
        ("Trade-weighted USD",  "Stronger $ → expensive gold abroad → less demand", PURPLE),
        ("Fed Balance Sheet",   "QE prints money → dollar dilution → gold up",     POS),
    ]
    col_w = Inches(2.4)
    gap = Inches(0.05)
    x0 = Inches(0.6)
    for i, (title, body, color) in enumerate(levers):
        x = x0 + (col_w + gap) * i
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), col_w, Inches(4.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CREAM
        card.line.color.rgb = color
        card.line.width = Pt(2)
        card.adjustments[0] = 0.06
        # number
        add_textbox(s, x + Inches(0.18), Inches(2.35), Inches(0.6), Inches(0.5),
                    f"0{i+1}", size=18, bold=True, color=color)
        # title
        add_textbox(s, x + Inches(0.18), Inches(2.85), col_w - Inches(0.36), Inches(0.9),
                    title, size=15, bold=True, color=INK)
        # body
        add_textbox(s, x + Inches(0.18), Inches(3.95), col_w - Inches(0.36), Inches(2.3),
                    body, size=11, color=INK_2)
    add_textbox(s, Inches(0.6), Inches(6.6), Inches(11), Inches(0.4),
                "Memorize these. They're 90% of gold analysis.",
                size=14, color=GOLD_300)

    # 5-9. THE 5 LEVERS — one slide each
    levers_detail = [
        # (number, title, current_value, what_it_is, why_it_matters, what_moves_it,
        #  effect_table_rows, strength, color)
        (
            "03", "Lever 1: US 10-Year Treasury Yield",
            "Currently ~4.2%. Forecast 4.07%.",
            "The interest rate the US government pays on 10-year bonds.",
            "This is gold's #1 competitor. If bonds pay 5% safely, why hold gold paying 0%?",
            ["Fed raises rates → 10y goes up",
             "Recession fears → flight to bonds → 10y goes down",
             "More government debt issuance → 10y goes up"],
            [
                ("Rises (4% → 5%)", "Gold falls", NEG),
                ("Falls (4% → 3%)", "Gold rises", POS),
            ],
            "STRONG · β = −0.094 · biggest lever",
            GOLD_500,
        ),
        (
            "04", "Lever 2: US Debt-to-GDP Growth",
            "Currently +0.5pp/yr. Forecast +1.3pp/yr.",
            "How fast US government debt grows relative to the economy. US debt ≈ 120% of GDP today.",
            "Massive debt → governments print money to pay it (inflation) or default. Gold hedges both.",
            ["Big stimulus / war / COVID → debt grows fast",
             "Tax hikes + spending cuts → debt grows slow",
             "Strong economy → GDP grows → debt/GDP falls"],
            [
                ("Rapidly rising (+3pp/yr)", "Gold rises", POS),
                ("Stable (~0)", "Neutral", INK_2),
                ("Falling", "Gold weak", NEG),
            ],
            "WEAK-MODERATE · β = +0.018 · slow-burn driver",
            BLUE,
        ),
        (
            "05", "Lever 3: US CPI Inflation",
            "Currently 4.5%. Forecast 2.8%.",
            "How fast prices for everyday stuff (food, gas, rent) rise in the US.",
            "Gold preserves purchasing power. The single cleanest 'inflation hedge' reason to own gold.",
            ["Money printing → more dollars chasing same goods → inflation up",
             "Supply chain breaks (war, energy crisis) → inflation up",
             "High interest rates eventually → inflation down"],
            [
                ("High (5%+)", "Gold strongly up", POS),
                ("Low (1-2%)", "Gold weak", NEG),
                ("Deflation", "Gold weak — cash wins", NEG),
            ],
            "VERY STRONG · β = +4.11 · biggest coefficient",
            NEG,
        ),
        (
            "06", "Lever 4: Trade-Weighted US Dollar",
            "Currently ~102.5. Forecast 100.",
            "Strength of the dollar vs major trading partners (EUR, JPY, GBP, CNY).",
            "Gold is priced in dollars. Strong $ → gold expensive abroad → less demand → gold (in USD) falls.",
            ["US rates higher than peers → capital flows in → DXY up",
             "Crisis (people flee to USD) → DXY up",
             "Fed printing → DXY down"],
            [
                ("Strong dollar (rising)", "Gold falls", NEG),
                ("Weak dollar (falling)", "Gold rises", POS),
            ],
            "STRONG · β = −0.87 · second biggest driver",
            PURPLE,
        ),
        (
            "07", "Lever 5: Federal Reserve Balance Sheet",
            "Currently ~$7T. Forecast $6.8T.",
            "Bonds + mortgages the Fed itself owns. QE = Fed buys (prints money). QT = Fed sells.",
            "When Fed prints money to buy bonds, more dollars in the system → dollar dilution → gold up. Usually.",
            ["Fed launches QE (crisis response) → balance sheet expands",
             "Fed runs QT (rolling off bonds) → balance sheet shrinks",
             "Long-term trend: from $0.9T (2007) → $9T peak (2022) → $7T now"],
            [
                ("Expanding (QE) — textbook", "Gold up", POS),
                ("Shrinking (QT) — textbook", "Gold weak", NEG),
                ("BUT 2022-2025: QT + gold rallied", "Sign flipped", NEG),
            ],
            "WEAK · β = −0.23 (negative — see note)",
            POS,
        ),
    ]
    for idx, (n, title, current, what, why, moves, table_rows, strength, color) in enumerate(levers_detail):
        s = prs.slides.add_slide(blank)
        add_bg(s, CREAM)
        # left side: accent bar in lever color
        accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SLIDE_H)
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()

        add_textbox(s, Inches(0.6), Inches(0.4), Inches(0.5), Inches(0.4),
                    n, size=18, bold=True, color=color)
        add_textbox(s, Inches(0.6), Inches(0.8), Inches(12), Inches(0.8),
                    title, size=26, bold=True, color=INK)
        add_textbox(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.4),
                    current, size=12, color=color, bold=True)

        # WHAT IT IS
        add_textbox(s, Inches(0.6), Inches(2.15), Inches(2.0), Inches(0.4),
                    "WHAT IT IS", size=10, bold=True, color=GOLD_700)
        add_textbox(s, Inches(2.7), Inches(2.15), Inches(10), Inches(0.7),
                    what, size=13, color=INK_2)

        # WHY IT MATTERS
        add_textbox(s, Inches(0.6), Inches(2.95), Inches(2.0), Inches(0.4),
                    "WHY IT MATTERS", size=10, bold=True, color=GOLD_700)
        add_textbox(s, Inches(2.7), Inches(2.95), Inches(10), Inches(0.7),
                    why, size=13, color=INK_2)

        # WHAT MOVES IT
        add_textbox(s, Inches(0.6), Inches(3.85), Inches(2.0), Inches(0.4),
                    "WHAT MOVES IT", size=10, bold=True, color=GOLD_700)
        add_bullet_list(s, Inches(2.7), Inches(3.85), Inches(10), Inches(1.3),
                        moves, size=11)

        # EFFECT ON GOLD
        add_textbox(s, Inches(0.6), Inches(5.25), Inches(2.0), Inches(0.4),
                    "EFFECT ON GOLD", size=10, bold=True, color=GOLD_700)
        add_table(s, Inches(2.7), Inches(5.25),
                  headers=["IF THIS HAPPENS", "GOLD DOES"],
                  rows=[(r[0], (r[1], r[2], True)) for r in table_rows],
                  col_widths=[Inches(5), Inches(5)],
                  row_height=Inches(0.4),
                  header_size=10, font_size=11)

        # STRENGTH chip
        chip_w = Inches(6.5)
        add_chip(s, Inches(0.6), Inches(6.8), strength,
                 color, CREAM, width=chip_w, height=Inches(0.4), size=10)

        add_footer(s, 4 + idx, slides_count_placeholder)

    # 10. THE RECIPE
    s = prs.slides.add_slide(blank)
    add_bg(s, GOLD_900)
    add_textbox(s, Inches(0.6), Inches(0.5), Inches(11), Inches(0.5),
                "08 · THE RECIPE", size=14, bold=True, color=GOLD_300)
    add_textbox(s, Inches(0.6), Inches(1.0), Inches(11), Inches(0.8),
                "The model, in one line",
                size=32, bold=True, color=CREAM)

    # The equation
    eq_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.6), Inches(2.3), Inches(12), Inches(1.4))
    eq_box.fill.solid()
    eq_box.fill.fore_color.rgb = CREAM
    eq_box.line.color.rgb = GOLD_500
    eq_box.line.width = Pt(2)
    eq_box.adjustments[0] = 0.05
    add_textbox(s, Inches(0.8), Inches(2.45), Inches(11.6), Inches(0.5),
                "Annual gold return =",
                size=14, color=INK_2, bold=True)
    add_textbox(s, Inches(0.8), Inches(2.85), Inches(11.6), Inches(0.9),
                "−1.6%  +  (−0.094 × Δ rates)  +  (0.018 × Δ debt)  +  (4.11 × inflation)  +  (−0.87 × Δ $)  +  (−0.23 × Δ Fed)",
                size=15, color=INK, font="Helvetica", bold=True)

    add_textbox(s, Inches(0.6), Inches(4.0), Inches(11), Inches(0.5),
                "In plain English:",
                size=14, bold=True, color=GOLD_300)
    add_bullet_list(s, Inches(0.6), Inches(4.5), Inches(12), Inches(2.5),
                    [
                        "Start with a baseline of −1.6% (the intercept).",
                        "Every 1 percentage-point increase in rates → subtract 0.094.",
                        "Every 1% inflation → ADD 4.11 (the biggest driver).",
                        "Every 1% dollar weakness → ADD 0.87.",
                        "Coefficients came from fitting 19 years of history (2007-2025).",
                        "R² = 0.61 → the model explains 61% of historical variance.",
                    ],
                    size=14, color=CREAM)

    # 11. WORKED EXAMPLE
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    add_left_accent(s)
    add_textbox(s, Inches(0.6), Inches(0.5), Inches(0.5), Inches(0.4),
                "09", size=18, bold=True, color=GOLD_500)
    add_textbox(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
                "Worked example · default forecast",
                size=28, bold=True, color=INK)
    add_textbox(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.5),
                "Plug the dashboard's default Qaurum-style macro path into our recipe:",
                size=13, color=INK_2)
    add_table(s, Inches(0.6), Inches(2.4),
              headers=["LEVER", "INPUT", "MATH", "CONTRIBUTION"],
              rows=[
                  ["US 10y rate", "4.07% (from 4.21%)", "−0.094 × (−0.028)", ("+0.003", POS, True)],
                  ["Debt/GDP",    "+1.3 pp/yr",         "0.018 × 1.3",       ("+0.023", POS, True)],
                  ["Inflation",   "2.8% per year",      "4.11 × 0.028",      ("+0.115", POS, True)],
                  ["USD",         "100 (from 102.5)",   "−0.87 × −0.005",    ("+0.004", POS, True)],
                  ["Fed BS",      "$6800B (from $7000B)", "−0.23 × −0.006",  ("+0.001", POS, True)],
                  ["Subtotal",    "",                   "",                  ("+0.146", GOLD_700, True)],
                  ["+ Intercept", "",                   "",                  ("−0.016", NEG, True)],
                  ["Predicted log-return", "",          "0.146 − 0.016",     ("+0.130", GOLD_700, True)],
              ],
              col_widths=[Inches(2.6), Inches(3.0), Inches(3.4), Inches(2.0)],
              header_size=11, font_size=11,
              row_height=Inches(0.4))

    add_textbox(s, Inches(0.6), Inches(6.0), Inches(12), Inches(0.4),
                "Predicted annual return = exp(0.130) − 1 ≈ +13.9%",
                size=18, bold=True, color=POS)
    add_textbox(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.5),
                "Driven almost entirely by inflation (+0.115 of the +0.146 sum). "
                "The chart will show inflation as the dominant contribution chip.",
                size=12, color=INK_2)
    add_footer(s, 10, slides_count_placeholder)

    # 12-15. 4 SCENARIOS
    scenarios = [
        (
            "10", "Scenario A: Soft Landing",
            "Mainstream view — Fed cuts, inflation cools, dollar weakens modestly",
            [["US 10y",        "3.50",  "rates fall"],
             ["Debt/GDP",      "1.0",   "modest growth"],
             ["Inflation",     "2.0",   "back to target"],
             ["USD",           "95",    "dollar weakens"],
             ["Fed BS",        "7000",  "flat"]],
            "+8% to +10%",
            "Bullish but not extreme. The default consensus call.",
            POS,
        ),
        (
            "11", "Scenario B: Stagflation",
            "1970s repeat — inflation sticks high, deficit explodes",
            [["US 10y",        "5.50",  "rates stay high"],
             ["Debt/GDP",      "2.5",   "deficit explodes"],
             ["Inflation",     "5.0",   "sticky"],
             ["USD",           "100",   "flat"],
             ["Fed BS",        "7500",  "flat"]],
            "+18% to +25%",
            "STRONG bull case. Historical 1970s: gold ran $35 → $850.",
            POS,
        ),
        (
            "12", "Scenario C: Goldilocks bear case for gold",
            "Rates rise, inflation crushed, USD strong, fiscal discipline",
            [["US 10y",        "6.00",  "rates rise"],
             ["Debt/GDP",      "0.0",   "fiscal discipline"],
             ["Inflation",     "1.5",   "low"],
             ["USD",           "115",   "USD strong"],
             ["Fed BS",        "5500",  "heavy QT"]],
            "−5% to flat",
            "Every lever against gold. Bearish.",
            NEG,
        ),
        (
            "13", "Scenario D: Crisis Year",
            "2008/2020 repeat — inflation spike, Fed prints massive QE, rates collapse",
            [["US 10y",        "2.50",  "rates collapse"],
             ["Debt/GDP",      "5.0",   "fiscal panic"],
             ["Inflation",     "6.0",   "inflation spike"],
             ["USD",           "88",    "USD weakens"],
             ["Fed BS",        "10000", "massive QE"]],
            "+25% to +35%",
            "Extreme bull case. Crisis = gold's home turf.",
            POS,
        ),
    ]
    for idx, (n, title, subtitle, rows, prediction, comment, color) in enumerate(scenarios):
        s = prs.slides.add_slide(blank)
        add_bg(s, CREAM)
        accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SLIDE_H)
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()
        add_textbox(s, Inches(0.6), Inches(0.5), Inches(0.5), Inches(0.4),
                    n, size=18, bold=True, color=color)
        add_textbox(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.7),
                    title, size=28, bold=True, color=INK)
        add_textbox(s, Inches(0.6), Inches(1.65), Inches(12), Inches(0.5),
                    subtitle, size=13, color=INK_2)

        add_textbox(s, Inches(0.6), Inches(2.4), Inches(5), Inches(0.4),
                    "SET THE INPUTS TO:", size=11, bold=True, color=GOLD_700)
        add_table(s, Inches(0.6), Inches(2.9),
                  headers=["LEVER", "VALUE", "MEANING"],
                  rows=rows,
                  col_widths=[Inches(2.0), Inches(1.5), Inches(2.7)],
                  row_height=Inches(0.4),
                  font_size=11)

        # Big prediction display
        pred_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(7.4), Inches(2.9), Inches(5.3), Inches(2.6))
        pred_box.fill.solid()
        pred_box.fill.fore_color.rgb = color
        pred_box.line.fill.background()
        pred_box.adjustments[0] = 0.05
        add_textbox(s, Inches(7.6), Inches(3.1), Inches(5), Inches(0.5),
                    "EXPECTED PREDICTION",
                    size=11, bold=True, color=CREAM)
        add_textbox(s, Inches(7.6), Inches(3.7), Inches(5), Inches(1.2),
                    prediction,
                    size=42, bold=True, color=CREAM)

        add_textbox(s, Inches(0.6), Inches(6.0), Inches(12), Inches(0.8),
                    comment, size=14, color=INK_2)
        add_footer(s, 11 + idx, slides_count_placeholder)

    # 16. HOW TO READ THE OUTPUT
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    add_left_accent(s)
    add_textbox(s, Inches(0.6), Inches(0.5), Inches(0.5), Inches(0.4),
                "14", size=18, bold=True, color=GOLD_500)
    add_textbox(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
                "How to read what you see",
                size=28, bold=True, color=INK)
    add_table(s, Inches(0.6), Inches(2.0),
              headers=["WHAT YOU SEE", "WHAT IT MEANS"],
              rows=[
                  ["R² = 0.61 badge", "Model explains 61% of historical gold-return variance. Anything above 0.5 is useful."],
                  ["Big '+13.9%' number", "Predicted gold return next year under your scenario inputs."],
                  ["Contribution chips", "How much each lever contributes. Biggest chip = your dominant driver."],
                  ["±1σ error bars", "Confidence range. ~2/3 of historical outcomes fell within these bars."],
                  ["β next to each input", "Sensitivity. Big β = lever moves the forecast a lot."],
                  ["Gold-highlighted cell", "You customised this from the default."],
                  ["'Reset N' button", "Wipe customisations, return to default."],
              ],
              col_widths=[Inches(3.5), Inches(8.5)],
              row_height=Inches(0.55),
              font_size=11)
    add_footer(s, 15, slides_count_placeholder)

    # 17. PRO WORKFLOW
    s = prs.slides.add_slide(blank)
    add_bg(s, GOLD_900)
    add_textbox(s, Inches(0.6), Inches(0.5), Inches(11), Inches(0.5),
                "15", size=14, bold=True, color=GOLD_300)
    add_textbox(s, Inches(0.6), Inches(1.0), Inches(11), Inches(0.8),
                "How a pro would actually use this",
                size=32, bold=True, color=CREAM)
    steps = [
        ("READ", "Spend a week with news + research. Form a macro view."),
        ("PLUG IN", "Enter your view into the inputs panel."),
        ("READ PREDICTION", "What return does the historical relationship imply?"),
        ("SANITY CHECK", "Are the contributions reasonable? Does the dominant lever match your gut?"),
        ("STRESS TEST", "What if you're wrong about inflation? Drop it 100bps. Does the call still hold?"),
        ("COMPARE", "Bloomberg analyst consensus, GS / JPM views. Where do you disagree?"),
        ("DECIDE", "Use this as ONE input among many — not gospel. Size accordingly."),
    ]
    y = Inches(2.3)
    for i, (head, body) in enumerate(steps):
        # number circle
        bullet = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(0.8), y + Inches(0.06), Inches(0.4), Inches(0.4))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = GOLD_500
        bullet.line.fill.background()
        add_textbox(s, Inches(0.8), y + Inches(0.06), Inches(0.4), Inches(0.4),
                    str(i + 1), size=14, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, Inches(1.5), y, Inches(2.0), Inches(0.55),
                    head, size=14, bold=True, color=GOLD_300)
        add_textbox(s, Inches(3.6), y, Inches(9.4), Inches(0.55),
                    body, size=13, color=CREAM)
        y += Inches(0.6)

    # 18. FINAL — MENTAL MODEL
    s = prs.slides.add_slide(blank)
    add_bg(s, CREAM)
    add_left_accent(s)
    add_textbox(s, Inches(0.6), Inches(0.5), Inches(0.5), Inches(0.4),
                "16", size=18, bold=True, color=GOLD_500)
    add_textbox(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
                "The mental model — keep this",
                size=30, bold=True, color=INK)

    # Two columns: UP, DOWN
    col_w = Inches(5.8)
    # UP
    up = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(0.6), Inches(2.0), col_w, Inches(4.5))
    up.fill.solid()
    up.fill.fore_color.rgb = RGBColor(0xE9, 0xF3, 0xE7)
    up.line.color.rgb = POS
    up.line.width = Pt(2)
    up.adjustments[0] = 0.05
    add_textbox(s, Inches(0.8), Inches(2.2), col_w - Inches(0.4), Inches(0.5),
                "Gold goes UP when:", size=18, bold=True, color=POS)
    add_bullet_list(s, Inches(0.8), Inches(2.9), col_w - Inches(0.4), Inches(3.5),
                    ["Inflation is high",
                     "Dollar is weak",
                     "Rates are low (especially real rates)",
                     "Government is in fiscal trouble",
                     "Crisis is brewing"],
                    size=14, color=INK_2, line_spacing=1.6)

    # DOWN
    dn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(6.7), Inches(2.0), col_w, Inches(4.5))
    dn.fill.solid()
    dn.fill.fore_color.rgb = RGBColor(0xF8, 0xE9, 0xE9)
    dn.line.color.rgb = NEG
    dn.line.width = Pt(2)
    dn.adjustments[0] = 0.05
    add_textbox(s, Inches(6.9), Inches(2.2), col_w - Inches(0.4), Inches(0.5),
                "Gold goes DOWN when:", size=18, bold=True, color=NEG)
    add_bullet_list(s, Inches(6.9), Inches(2.9), col_w - Inches(0.4), Inches(3.5),
                    ["Inflation is killed",
                     "Dollar is strong",
                     "Rates are high (especially real rates)",
                     "Fiscal house in order",
                     "Everything is calm"],
                    size=14, color=INK_2, line_spacing=1.6)

    add_textbox(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
                "You're now equipped to use the Forecast tab as a real gold analyst. Have fun.",
                size=13, color=INK_2, align=PP_ALIGN.CENTER)
    add_footer(s, 16, slides_count_placeholder)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    size_kb = OUT_PATH.stat().st_size / 1024
    print(f"[ppt] wrote {OUT_PATH.relative_to(ROOT)} ({size_kb:.0f} KB, {len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
